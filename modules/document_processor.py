"""
Module: document_processor.py
Phase: 2 - Document Processing
Responsibility: Extract raw text from PDF, PPTX, TXT files ONLY.
                No chunking, no embeddings — just clean text.
"""

import os
from pathlib import Path


def extract_text(file_path: str) -> str:
    """
    Main entry point. Detects file type and routes to correct extractor.

    Args:
        file_path: Absolute or relative path to the uploaded file.

    Returns:
        Cleaned extracted text as a single string.

    Raises:
        ValueError: If file type is unsupported.
        FileNotFoundError: If file does not exist.
    """
    path = Path(file_path)

    if not path.exists():
        raise FileNotFoundError(f"File not found: {file_path}")

    ext = path.suffix.lower()

    extractors = {
        ".pdf": _extract_from_pdf,
        ".pptx": _extract_from_pptx,
        ".txt": _extract_from_txt,
        ".md": _extract_from_txt,
    }

    if ext not in extractors:
        raise ValueError(
            f"Unsupported file type: '{ext}'. "
            f"Supported: {list(extractors.keys())}"
        )

    raw_text = extractors[ext](str(path))
    return _clean_text(raw_text)


# ──────────────────────────────────────────────
# Private extractors (one per file type)
# ──────────────────────────────────────────────

def _extract_from_pdf(file_path: str) -> str:
    """Extract text from a PDF using pypdf."""
    try:
        from pypdf import PdfReader
    except ImportError:
        raise ImportError("Run: pip install pypdf")

    reader = PdfReader(file_path)
    pages_text = []

    for i, page in enumerate(reader.pages):
        text = page.extract_text()
        if text:
            pages_text.append(f"[Page {i + 1}]\n{text}")

    if not pages_text:
        raise ValueError(
            "No text could be extracted from the PDF. "
            "It may be a scanned image PDF — try converting to text first."
        )

    return "\n\n".join(pages_text)


def _extract_from_pptx(file_path: str) -> str:
    """Extract text from PowerPoint slides using python-pptx."""
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("Run: pip install python-pptx")

    prs = Presentation(file_path)
    slides_text = []

    for i, slide in enumerate(prs.slides):
        slide_content = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        slide_content.append(line)
        if slide_content:
            slides_text.append(f"[Slide {i + 1}]\n" + "\n".join(slide_content))

    return "\n\n".join(slides_text)


def _extract_from_txt(file_path: str) -> str:
    """Read a plain text or markdown file."""
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def _clean_text(text: str) -> str:
    """
    Light cleaning: remove excessive whitespace, null bytes, etc.
    Does NOT remove newlines — those help with chunking later.
    """
    # Remove null bytes
    text = text.replace("\x00", "")
    # Collapse 3+ consecutive newlines into 2
    import re
    text = re.sub(r"\n{3,}", "\n\n", text)
    # Strip leading/trailing whitespace
    return text.strip()


# ──────────────────────────────────────────────
# Quick test (run this file directly)
# ──────────────────────────────────────────────
if __name__ == "__main__":
    import sys

    if len(sys.argv) < 2:
        print("Usage: python document_processor.py <path_to_file>")
        sys.exit(1)

    test_path = sys.argv[1]
    print(f"Extracting text from: {test_path}")
    text = extract_text(test_path)
    print(f"\n✅ Extracted {len(text)} characters")
    print("\n--- Preview (first 500 chars) ---")
    print(text[:500])
